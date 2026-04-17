from pptx import Presentation

prs = Presentation()

# Slide 1: Title
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Secure, Local Agent Architecture"
slide.placeholders[1].text = "Layered Guardrails for MCP-Enabled LLMs in Air-Gapped Environments"

# Slide 2: Architectural Philosophy
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Architectural Philosophy"
tf = slide.placeholders[1].text_frame
tf.text = "Objective: Implement enterprise-grade security for local LLM agents without relying on third-party APIs."
p = tf.add_paragraph(); p.text = "Decoupled Design: Guardrails are isolated in a standalone Python class."; p.level = 1
p = tf.add_paragraph(); p.text = "Hardware Optimized: Prioritizes deterministic, CPU-efficient libraries to reserve memory for the main model (qwen3.5:9b)."; p.level = 1
p = tf.add_paragraph(); p.text = "Privacy First: 100% local execution ensures sensitive data never leaves the 8-core infrastructure."; p.level = 1
p = tf.add_paragraph(); p.text = "Native Integration: Leverages modern LlamaIndex FunctionAgent and Memory modules."; p.level = 1

# Slide 3: Layer 1
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Layer 1: Input Guardrails"
tf = slide.placeholders[1].text_frame
tf.text = "Goal: Sanitize and reject malicious inputs before allocating inference resources."
p = tf.add_paragraph(); p.text = "Prompt Injection Blocking: Uses lightweight regex to catch attempts to override system prompts."; p.level = 1
p = tf.add_paragraph(); p.text = "Deterministic Filtering: Employs better_profanity for instant, hash-map-based rejection."; p.level = 1
p = tf.add_paragraph(); p.text = "Zero-Shot PII Redaction: Integrates Microsoft Presidio (via spaCy) to anonymize sensitive data into safe tags."; p.level = 1

# Slide 4: Layer 2
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Layer 2: Structural Guardrails"
tf = slide.placeholders[1].text_frame
tf.text = "Goal: Guarantee the LLM interacts safely and accurately with connected MCP servers."
p = tf.add_paragraph(); p.text = "Strict Schema Enforcement: Uses Pydantic to define exact input types and constraints."; p.level = 1
p = tf.add_paragraph(); p.text = "Self-Correcting Tool Usage: Pydantic validation errors prompt the FunctionAgent to auto-correct."; p.level = 1
p = tf.add_paragraph(); p.text = "Immutable System Prompt: Establishes rigid boundaries at the agent instantiation layer."; p.level = 1

# Slide 5: Layer 3
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Layer 3: Output Guardrails"
tf = slide.placeholders[1].text_frame
tf.text = "Goal: Ensure the final response is safe, professional, and free of internal system leaks."
p = tf.add_paragraph(); p.text = "Secondary PII & Profanity Scan: Runs text back through Presidio and better_profanity."; p.level = 1
p = tf.add_paragraph(); p.text = "Asymmetric 'LLM-as-a-Judge': Offloads semantic validation to an ultra-lightweight model (qwen3.5:0.5b)."; p.level = 1
p = tf.add_paragraph(); p.text = "Leak Prevention: Rapidly scans output for exposed Python tracebacks, raw code, or leaked instructions."; p.level = 1

# Slide 6: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Why This Beats 'Off-the-Shelf'"
tf = slide.placeholders[1].text_frame
tf.text = "Summary of the Custom Guardrail Advantage:"
p = tf.add_paragraph(); p.text = "Zero Framework Bloat: Eliminates heavy dependencies (like guardrails-ai)."; p.level = 1
p = tf.add_paragraph(); p.text = "Maximum CPU Efficiency: Offloads structural checks to Pydantic and text checks to spaCy/regex."; p.level = 1
p = tf.add_paragraph(); p.text = "Total Control: Custom error handling mapped directly to API HTTP status codes (400, 403, 422)."; p.level = 1

# Save the presentation
output_path = 'Secure_Agent_Architecture.pptx'
prs.save(output_path)
print(f"Presentation saved to {output_path}")
